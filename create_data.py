"""
create_data.py
--------------
Ingestion pipeline for the FinRAG backend.

ZERO langchain dependencies in the loaders themselves — uses only the
Python stdlib (plus light, swappable third-party parsers per file type).
The embedding + vector store object is built once in main.py and passed
in here, exactly like the original CLI version of this project.

Supported file types (matches the 8 cards on the upload page):
  1. pdf      .pdf
  2. excel    .xlsx / .xls
  3. csv      .csv
  4. json     .json
  5. sqlite   .db / .sqlite / .sqlite3
  6. ppt      .pptx / .ppt
  7. html     .html / .htm
  8. text     .txt
"""

import csv
import io
import json
import sqlite3
from html.parser import HTMLParser
from pathlib import Path

# ── File-type registry (ids match js/upload.js FILE_TYPES) ──────────────────

FILE_TYPES = {
    "pdf":    {"label": "PDF",        "ext": [".pdf"]},
    "excel":  {"label": "Excel",      "ext": [".xlsx", ".xls"]},
    "csv":    {"label": "CSV",        "ext": [".csv"]},
    "json":   {"label": "JSON",       "ext": [".json"]},
    "sqlite": {"label": "SQLite",     "ext": [".db", ".sqlite", ".sqlite3"]},
    "ppt":    {"label": "PowerPoint", "ext": [".pptx", ".ppt"]},
    "html":   {"label": "HTML",       "ext": [".html", ".htm"]},
    "text":   {"label": "Text",       "ext": [".txt"]},
}


class IngestError(Exception):
    """Raised for any recoverable ingestion problem (bad file, bad type, etc.)."""


def validate_extension(filename: str, file_type: str) -> None:
    """Raise IngestError if `filename`'s extension doesn't match `file_type`."""
    if file_type not in FILE_TYPES:
        raise IngestError(f"Unsupported file type '{file_type}'.")
    suffix = Path(filename).suffix.lower()
    allowed = FILE_TYPES[file_type]["ext"]
    if suffix not in allowed:
        raise IngestError(
            f"'{filename}' does not match the selected type "
            f"({FILE_TYPES[file_type]['label']}). Expected: {', '.join(allowed)}"
        )


# ── Pure-Python text chunker (no langchain needed) ───────────────────────────

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 150):
    """Split text into overlapping chunks of ~chunk_size characters."""
    chunks = []
    start = 0
    length = len(text)
    if length == 0:
        return chunks
    while start < length:
        end = min(start + chunk_size, length)
        piece = text[start:end].strip()
        if piece:
            chunks.append(piece)
        if end == length:
            break
        start = end - overlap
    return chunks


def _group_rows_into_blocks(rows, max_len: int = 1000):
    """Pack a list of short text rows into ~max_len character blocks."""
    block, blocks, cur_len = [], [], 0
    for row in rows:
        if cur_len + len(row) > max_len and block:
            blocks.append("\n".join(block))
            block, cur_len = [], 0
        block.append(row)
        cur_len += len(row)
    if block:
        blocks.append("\n".join(block))
    return blocks


# ── File loaders — each takes raw `bytes` and returns a list[str] of sections ─

def load_pdf(data: bytes):
    try:
        import pypdf
    except ImportError as e:
        raise IngestError("pypdf is not installed on the server.") from e
    texts = []
    reader = pypdf.PdfReader(io.BytesIO(data))
    for page in reader.pages:
        t = page.extract_text() or ""
        if t.strip():
            texts.append(t)
    return texts


def load_text(data: bytes):
    return [data.decode("utf-8", errors="replace")]


def load_csv(data: bytes):
    text = data.decode("utf-8-sig", errors="replace")
    rows = []
    reader = csv.DictReader(io.StringIO(text))
    for row in reader:
        line = " | ".join(f"{k}: {v}" for k, v in row.items() if v)
        if line.strip():
            rows.append(line)
    return _group_rows_into_blocks(rows)


def load_excel(data: bytes):
    try:
        import openpyxl
    except ImportError as e:
        raise IngestError("openpyxl is not installed on the server.") from e
    rows = []
    wb = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        all_rows = list(ws.iter_rows(values_only=True))
        if not all_rows:
            continue
        headers = [str(h) if h is not None else f"col_{i}" for i, h in enumerate(all_rows[0])]
        for row in all_rows[1:]:
            pairs = [
                f"{headers[i]}: {row[i]}"
                for i in range(min(len(headers), len(row)))
                if row[i] is not None
            ]
            if pairs:
                rows.append(f"[{sheet_name}] " + " | ".join(pairs))
    return _group_rows_into_blocks(rows)


def load_json(data: bytes):
    text = data.decode("utf-8", errors="replace")
    try:
        obj = json.loads(text)
    except json.JSONDecodeError as e:
        raise IngestError(f"Invalid JSON file: {e}") from e

    lines = []

    def walk(node, path):
        if isinstance(node, dict):
            for k, v in node.items():
                walk(v, f"{path}.{k}" if path else str(k))
        elif isinstance(node, list):
            for i, v in enumerate(node):
                walk(v, f"{path}[{i}]")
        else:
            lines.append(f"{path}: {node}")

    walk(obj, "")
    return _group_rows_into_blocks(lines)


def load_sqlite(data: bytes, tmp_dir: Path):
    # sqlite3 needs a real file path — write to a temp file in the caller's dir.
    tmp_path = tmp_dir / "_ingest_tmp.sqlite"
    tmp_path.write_bytes(data)
    rows = []
    try:
        conn = sqlite3.connect(str(tmp_path))
        cur = conn.cursor()
        cur.execute("SELECT name FROM sqlite_master WHERE type='table';")
        tables = [t[0] for t in cur.fetchall()]
        for table in tables:
            cur.execute(f"SELECT * FROM [{table}];")
            cols = [d[0] for d in cur.description]
            for row in cur.fetchall():
                line = f"[{table}] " + " | ".join(
                    f"{c}: {v}" for c, v in zip(cols, row) if v is not None
                )
                rows.append(line)
        conn.close()
    finally:
        tmp_path.unlink(missing_ok=True)
    return _group_rows_into_blocks(rows)


def load_pptx(data: bytes):
    try:
        from pptx import Presentation
    except ImportError as e:
        raise IngestError("python-pptx is not installed on the server.") from e
    prs = Presentation(io.BytesIO(data))
    texts = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for r in shape.table.rows:
                    parts.append(" | ".join(c.text for c in r.cells))
        if parts:
            texts.append(f"[Slide {i}]\n" + "\n".join(parts))
    return texts


class _TextOnlyHTMLParser(HTMLParser):
    """Minimal stdlib HTML -> text extractor (skips script/style tags)."""

    def __init__(self):
        super().__init__()
        self._skip = False
        self.parts = []

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style"):
            self._skip = True

    def handle_endtag(self, tag):
        if tag in ("script", "style"):
            self._skip = False

    def handle_data(self, data):
        if not self._skip:
            stripped = data.strip()
            if stripped:
                self.parts.append(stripped)


def load_html(data: bytes):
    text = data.decode("utf-8", errors="replace")
    parser = _TextOnlyHTMLParser()
    parser.feed(text)
    full_text = "\n".join(parser.parts)
    return [full_text] if full_text.strip() else []


_LOADERS = {
    "pdf": load_pdf,
    "excel": load_excel,
    "csv": load_csv,
    "json": load_json,
    "ppt": load_pptx,
    "html": load_html,
    "text": load_text,
    # "sqlite" handled specially below (needs a temp directory)
}


# ── Main entry point (called from main.py's /api/upload route) ──────────────

def ingest_file(filename: str, data: bytes, file_type: str, vector_store, tmp_dir: Path):
    """
    Full pipeline: validate -> load -> chunk -> embed/store.

    vector_store: a pre-built Chroma instance (one per user/session),
                  passed in from main.py so this module never has to
                  import langchain itself.
    tmp_dir:      a writable scratch directory (only used by the sqlite loader).

    Returns: number of chunks stored.
    Raises:  IngestError on any validation/parsing problem.
    """
    validate_extension(filename, file_type)

    if not data:
        raise IngestError("Uploaded file is empty.")

    if file_type == "sqlite":
        raw_sections = load_sqlite(data, tmp_dir)
    else:
        raw_sections = _LOADERS[file_type](data)

    if not raw_sections:
        raise IngestError("No readable content could be extracted from the file.")

    chunks = []
    for section in raw_sections:
        chunks.extend(chunk_text(section))

    if not chunks:
        raise IngestError("File content was extracted but produced no usable chunks.")

    metadatas = [{"source": filename, "chunk": i} for i in range(len(chunks))]
    vector_store.add_texts(chunks, metadatas=metadatas)
    return len(chunks)
